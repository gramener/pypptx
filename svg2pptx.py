"""
Converts an SVG shape into a Microsoft Office object, and saves as .pptx

Source: https://github.com/gramener/pypptx

Usage: python svg2pptx.py filename.svg
"""
import re
from lxml import etree, html
from lxml.builder import ElementMaker
from pypptx import a, p, shape, color, nsmap, cust_shape, cust_table
from color import rgba


re_ns = re.compile(r'({.*?})?(.*)')
re_path = re.compile(r'[mMzZlLhHvVcCsSqQtTaA]|[\+\-]?[\d\.e]+')

def interpret_str(val):
    if val:
        rn = re.compile(r'([\d\.-]+)')
        match = rn.search(val)
        value = match.group(1)
        if val.startswith('-'):
            val = str(0)
        elif val.endswith('%'):
            val = str(int(value) * 0.12)
        elif val.endswith('em'):
            val = str(value * 10 + 8)
        elif val.endswith('pt'):
            val = str(int(value + 6))
        else:
            val = value
    return val


def msclr(color):
    r, g, b, a = rgba(color)
    return '%02x%02x%02x' % (255*r, 255*g, 255*b)

def css_style(style):
    e = {}
    attrs = [x for x in style.split(";") if x != '']

    for attr in attrs:
        keys, values = attr.split(':')
        key, value = keys.split(), values.split()
        e.update(dict(zip(key, value)))
    return e


def tag_attrs(keys, values, e):
    parent = e.getparent()
    attrs_dict = dict(zip(keys, values))
    if 'style' in keys:
        del attrs_dict['style']
        style_dict = css_style(e.get('style'))
        attrs_dict.update(style_dict)
    if parent.tag == 'g':
        g_keys = parent.keys()
        values = parent.values()
        g_dict = dict(zip(g_keys, values))
        attrs_dict.update(g_dict)
        if 'style' in g_keys:
            del g_dict['style']
            g_attrs_dict = css_style(parent.get('style'))
            attrs_dict.update(g_attrs_dict)
    return attrs_dict

def translate(e):
    gtag = e.xpath('ancestor::*/@transform')
    atag = e.get('transform')
    if atag is not None and atag.startswith('translate'):
        xy = re.findall('\d*\.?\d+', atag)
        xy2 = re.findall('\d*\.?\d+', gtag[0] if len(gtag) >= 1 else '0,0')
        x, y = str(float(xy[0]) + float(xy2[0])), str(float(xy[1]) + float(xy2[1]))
    elif gtag:
        xy = re.findall('\d*\.?\d+', gtag[-1])
        xy2 = re.findall('\d*\.?\d+', gtag[0] if len(gtag) > 1 else '0,0')
        x, y = str(float(xy[0]) + float(xy2[0])), str(float(xy[1]) + float(xy2[1]))
    else:
        x, y = 0, 0
    return x, y



class Draw(object):
    def __init__(self, slide, width, height):
        self.slide = slide
        self.shapes = slide._element.find('.//p:spTree', namespaces=nsmap)
        # TODO: Replace 9999... with slide width and height
        self.x = lambda x: int(float(x) * 9999999 / width)
        self.y = lambda y: int(float(y) * 7777777 / height)


    def _shape_attrs(function):
        def wrapped(self, e):
            shape = function(self, e)
            # TODO: tooltip
            # child = [x.tag for x in e.getchildren()]
            # title_text = [x.text for x in e.getchildren()]
            # if 'title' in child:
            #     shape.find('.//a:cNvPr', namespaces=nsmap).append(
            #         a.hlinkClick, action="ppaction://hlinksldjump", tooltip=title_text)
            #     print title_text
            tag = function.__name__
            keys = e.keys()
            values = e.values()

            def styles(keys):
                def clr_grad(color):
                    if color.startswith('rgba('):
                        r, g, b, a = rgba(color)
                        return  '%d' % int(a*100000)
                    elif 'opacity' in keys:
                        return '%d' % int(float(e.get('opacity')) * 100000)
                    else:
                        return '%d' % 100000


                # TODO: Optimize
                if 'fill' in keys:
                    if e.get('fill') == 'none':
                        shape.spPr.append(a.noFill())
                    else:
                        shape.spPr.append(a.solidFill(a.srgbClr(a.alpha(val=str(clr_grad(e.get('fill')))),
                            val=str(msclr(e.get('fill'))))))
                elif not 'fill' in keys:
                    if tag not in ['line']:
                        shape.spPr.append(a.solidFill(color(srgbClr='000000')))

                if 'stroke' in keys and 'stroke-width' in keys:
                    shape.spPr.append(a.ln(a.solidFill(a.srgbClr(a.alpha(val=str(clr_grad(e.get('stroke')))),
                        val=str(msclr(e.get('stroke'))))),
                        w=str(int(float(interpret_str(e.get('stroke-width')))*12700))))
                elif 'stroke' in keys:
                    if e.get('stroke') == 'none':
                        shape.spPr.append(a.ln(a.noFill()))
                    else:
                        shape.spPr.append(a.ln(a.solidFill(a.srgbClr(a.alpha(val=str(clr_grad(e.get('stroke')))), val=str(msclr(e.get('stroke')))))))
                elif 'stroke' and 'fill' not in keys:
                    shape.spPr.append(a.ln(a.solidFill(color(srgbClr='000000'))))
                elif not 'stroke' and 'fill' in keys:
                    if tag in ['rect']:
                        shape.spPr.append(a.ln(a.noFill()))
                    elif tag in ['circle','ellipse']:
                        shape.spPr.append(a.ln(a.solidFill(color(srgbClr=msclr(e.get('fill'))))))
                    elif tag in ['path', 'line']:
                        shape.spPr.append(a.ln(a.solidFill(color(srgbClr='000000'))))
                return shape


            e = tag_attrs(keys, values, e)
            keys = e.keys()
            styles(keys)

            return shape
        return wrapped

    @_shape_attrs
    def circle(self, e):
        ax, ay = translate(e)
        x = float(e.get('cx', 0)) + float(ax)
        y = float(e.get('cy', 0)) + float(ay)
        r = float(e.get('r', 0))
        shp = shape('ellipse', self.x(x - r), self.y(y - r),
                               self.x(2 * r), self.y(2 * r))
        self.shapes.append(shp)
        return shp

    @_shape_attrs
    def ellipse(self, e):
        ax, ay = translate(e)
        x = float(e.get('cx', 0)) + float(ax)
        y = float(e.get('cy', 0)) + float(ay)
        rx = float(e.get('rx', 0))
        ry = float(e.get('ry', 0))
        shp = shape('ellipse', self.x(x - rx), self.y(y - ry),
                               self.x(2 * rx), self.y(2 * ry))
        self.shapes.append(shp)
        return shp

    @_shape_attrs
    def rect(self, e):
        ax, ay = translate(e)
        x = float(interpret_str(e.get('x', 0))) + float(ax)
        y = float(interpret_str(e.get('y', 0))) + float(ay)
        keys = e.keys()
        shp_name = 'roundRect' if 'rx' in keys and 'ry' in keys else 'rect'
        shp = shape(shp_name,
            self.x(x),
            self.y(y),
            self.x(interpret_str(e.get('width', 0))),
            self.y(interpret_str(e.get('height', 0)))
        )
        self.shapes.append(shp)
        return shp

    @_shape_attrs
    def line(self, e):
        ax, ay = translate(e)
        x1 = self.x(float(interpret_str(e.get('x1', 0))) + float(ax))
        y1 = self.y(float(interpret_str(e.get('y1', 0))) + float(ay))
        x2 = self.x(float(interpret_str(e.get('x2', 0))) + float(ax))
        y2 = self.y(float(interpret_str(e.get('y2', 0))) + float(ay))
        ax1 = x1 if x2 > x1 else x2
        ax2 = x2 if x1 < x2 else x1
        ay1 = y1 if y2 > y1 else y2
        ay2 = y2 if y1 < y2 else y1
        shp = shape('line', ax1, ay1, ax2-ax1, ay2-ay1)
        self.shapes.append(shp)
        return shp

    def text(self, e):
        keys = e.keys()
        values = e.values()
        txt = e.text
        def txt_anchor():
            anchor_dict = {'hanging':'t', 'middle':'ctr', True:'t', False:'ctr', 'left':'ctr'}
            if 'dominant-baseline' in keys:
                anchor = anchor_dict[e.get('dominant-baseline')]
            elif 'dy' in keys:
                em = float(re.findall(".\d+", e.get('dy'))[0]) > 0.5
                anchor = anchor_dict[em]
            # elif 'text-anchor' in keys:
            #     anchor = anchor_dict[e.get('text-anchor')]
            else:
                anchor = 'ctr'
            return anchor

        def txt_align():
            if 'text-anchor' in keys:
                align_dict = {'end':'r', 'middle':'ctr', 'start':'l', 'left':'l'}
                align = align_dict[e.get('text-anchor')]
            else:
                align = 'l'
            return align

        if not e.text:
            return
        ax, ay = translate(e)
        x = float(interpret_str(e.get('x', 0))) + float(ax)
        y = float(interpret_str(e.get('y', 0))) + float(ay)
        shp = shape('rect', self.x(x), self.y(y), self.x(0), self.y(0))
        if 'transform' in keys:
            t_key = e.get('transform')
            rotate = str(int(t_key[(t_key.find('rotate')+7):-1].split()[0])*60000)
            shp.find('.//a:xfrm', namespaces=nsmap).set('rot', rotate)
        def text_style(keys, txt):

            bold = '1' if 'font-weight' in keys else '0'

            fill_text_ml = a.solidFill(color(srgbClr=msclr(e.get('fill') if 'fill' in keys else 'black')))

            autofit_ml = a.normAutofit(fontScale="62500", lnSpcReduction="20000")    # Auto fit

            font_size = str(int(float(interpret_str(e.get('font-size')))*100)) if 'font-size' in keys else '1600'


            shp.append(p.txBody(a.bodyPr(anchor=txt_anchor(), wrap='none'),
            a.p(a.pPr(algn=txt_align()), a.r(a.rPr(fill_text_ml, lang='en-US', sz=font_size, b=bold, dirty='0', smtClean='0'),
                    a.t(txt)))))

            return shp

        e = tag_attrs(keys, values, e)
        keys = e.keys()
        text_style(keys, txt)

        self.shapes.append(shp)
        return shp

    # def table(self, e):
    #     thead_th = e.xpath('//table//th')
    #     tbody_td = e.xpath('//table//td')
    #     x = e.xpath('//table//tr')
    #     for y in x:
    #         for z in y:
    #             print z.text
    #     rows = len(thead_th)

    #     # cust_table(x, y, cx, cy)
    #     shp = cust_table('464016', '1397000', '8188664', '1982034' )
    #     gridcol = []
    #     th_list = []
    #     td_list = []

    #     for th in thead_th:
    #         gc = a.gridCol(w="744424")
    #         gridcol.append(gc)
    #         th_list.append(th.text)
    #     for td in tbody_td:
    #         td_list.append(td.text)

    #     texts = th_list+td_list
    #     text_values = [texts[i:i+rows] for i in range(0, len(texts), rows)]

    #     shp.find('.//a:tbl', namespaces=nsmap).append(a.tblGrid(*gridcol))
    #     for row in text_values:
    #         shp.find('.//a:tbl', namespaces=nsmap).append(a.tr(h='841233'))
    #         for val in row:
    #             print val


    #     self.shapes.append(shp)
    #     return shp


    @_shape_attrs
    def path(self, e):
        pathstr = re_path.findall(e.get('d', '')) if 'nan' not in e.get('d') else []
        n, length, cmd, relative, shp = 0, len(pathstr), None, False, None
        x1, y1 = 0, 0
        ax, ay = translate(e)
        xy = lambda n: (float(pathstr[n]) + (x1 if relative else 0) + float(ax),
                        float(pathstr[n + 1]) + (y1 if relative else 0) + float(ay))

        shp = cust_shape(x1, y1, self.x(100000), self.y(100000))
        path = a.path(w=str(self.x(100000)), h=str(self.y(100000)))
        shp.find('.//a:custGeom', namespaces=nsmap).append(
            a.pathLst(path))

        while n < length:
            if pathstr[n].lower() in 'mzlhvcsqta':
                cmd = pathstr[n].lower()
                relative = str.islower(pathstr[n])
                n += 1


            if cmd == 'm':
                x1, y1 = xy(n)
                path.append(a.moveTo(a.pt(x=str(self.x(x1)), y=str(self.y(y1)))))
                n += 2


            elif cmd == 'z':
                path.append(a.close())

            elif cmd == 'l':
                x1, y1 = xy(n)
                path.append(a.lnTo(a.pt(x=str(self.x(x1)), y=str(self.y(y1)))))
                n += 2


            elif cmd == 'c':
                xc1, yc1 = xy(n)
                xc2, yc2 = xy(n + 2)
                x1, y1 = xy(n + 4)
                path.append(a.cubicBezTo(a.pt(x=str(self.x(xc1)), y=str(self.y(yc1))),
                    a.pt(x=str(self.x(xc2)), y=str(self.y(yc2))),
                    a.pt(x=str(self.x(x1)), y=str(self.y(y1)))))
                n += 6


            #TODO blockArc:
            #elif cmd == 'a':
            #    x1, y1 = xy(n)
            #    cx, cy = xy(n + 5)
            #    shp = shape('blockArc', self.x(x1), self.y(y1), self.x(cx), self.y(cy))
            #    n += 7

            elif cmd == 'a':
                wR, hR = xy(n)
                stAng, swAng = xy(n + 5)
                path.append(a.arcTo(
                    wR=str(self.x(wR)), hR=str(self.y(hR)),
                    stAng=str(self.x(stAng)), swAng=str(self.y(swAng))))
                n += 7

        self.shapes.append(shp)
        return shp



def svg2mso(slide, svg, width=940, height=None):
    if width is not None and height is None:
        height = width * 3 / 4
    elif width is None and height is not None:
        width = height * 4 / 3

    # Convert tree into an lxml etree if it's not one
    if not hasattr(svg, 'iter'):
        svg = etree.parse(svg) if hasattr(svg, 'read') else etree.fromstring(svg)

    # Take all the tags and draw it
    draw = Draw(slide, width, height)
    valid_tags = set(tag for tag in dir(draw) if not tag.startswith('_'))
    for e in svg.iter(tag=etree.Element):
        match = re_ns.match(e.tag)
        if not match:
            continue

        tag = match.groups()[-1]
        if tag in valid_tags:
            getattr(draw, tag)(e)


if __name__ == '__main__':
    import argparse

    parser = argparse.ArgumentParser(description=__doc__.strip())
    parser.add_argument('--layout',
        default='layout15x12.pptx',
        help='PPTX file to use to create blank slide')
    parser.add_argument('--output',
        default='output.pptx',
        help='Output PPTX file name')
    parser.add_argument('svgfile')
    args = parser.parse_args()

    from pptx import Presentation

    ppt = Presentation(args.layout)
    blank_slidelayout = ppt.slidelayouts[6]
    slide = ppt.slides.add_slide(blank_slidelayout)
    tree = html.parse(open(args.svgfile))
    svg2mso(slide, tree)
    ppt.save(args.output)
