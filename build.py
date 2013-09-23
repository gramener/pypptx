"""
Creates a PPTX with all SVG files in the tests/ folder.
Useful for testing.

Usage: python build.py
"""

if __name__ == '__main__':
    import glob
    from lxml import etree
    from svg2pptx import svg2mso
    from pptx import Presentation

    Presentation = Presentation()
    blank_slidelayout = Presentation.slidelayouts[6]

    for svgfile in glob.glob('tests/*.svg'):
        print svgfile
        slide = Presentation.slides.add_slide(blank_slidelayout)
        tree = etree.parse(open(svgfile))
        title = etree.SubElement(tree.getroot(), 'text', x="300", y="20")
        title.text = svgfile
        try:
            svg2mso(slide, tree)
        except:
            print 'Failed'
            pass

    Presentation.save("test.pptx")
