from pptx import Presentation
from pypptx import nsmap, a, p, shape, color

prs = Presentation()
slide = prs.slides.add_slide(prs.slidelayouts[6])
shapes = slide._element.find('.//p:spTree', namespaces=nsmap)

shp = shape('ellipse', 0, 0, 999999, 999999)
shapes.append(shp)

# Fill with a scheme colour
shp.spPr.append(a.solidFill(color(
    schemeClr='accent2',  # 2nd theme colour
    tint='50%',           # 50% white mixed
    alpha='30%'           # 30% opaque, 70% transparent
)))

shp = shape('ellipse', 999999, 0, 999999, 999999)
shapes.append(shp)

# Fill with an RGB colour
shp.spPr.append(a.solidFill(color(
    srgbClr='FF0000',     # Red
    shade='50%',          # 50% black mixed
    sat='30%'             # 30% saturation
)))

shp = shape('ellipse', 0, 999999, 999999, 999999)
shapes.append(shp)

# Fill with an RGB colour
shp.spPr.append(a.gradFill(
    a.gsLst(
        a.gs(color(schemeClr='accent2', tint= '0%'), pos="0"),
        a.gs(color(schemeClr='accent2', tint='20%'), pos="50000"),
        a.gs(color(schemeClr='accent2', tint='40%'), pos="100000"),
    ),
    a.lin(ang='2700000', scaled='1'), # out of 21600000 = 1/8 = 45 degrees
))

# Add a line
shp.spPr.append(a.ln(
    a.solidFill(color(        # Solid fill with
        schemeClr='accent2',  # 2nd theme colour
        shade='20%',          # 20% black mixed
        alpha='50%',          # 50% transparent
    )),
    w='3175',                 # 0.25pt stroke width
))

# Add text
shp.append(p.txBody(
    a.bodyPr(anchor='ctr'),   # vertically center the text
    a.p(
        a.pPr(algn='ctr'),     # horizontally center the text
        a.r(a.t('abc')),
)))
prs.save('sample.pptx')
